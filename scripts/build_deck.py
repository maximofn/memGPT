"""Genera charla-memgpt.pptx a partir del paper, el resumen y la implementación.

Ejecutar:
    conda deactivate && source .venv/bin/activate
    uv run scripts/build_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "charla-memgpt.pptx"

# ---------- Paleta (60-30-10) ----------
BG = RGBColor(0x0E, 0x14, 0x22)           # azul nocturno (dominante)
SURFACE = RGBColor(0x16, 0x1F, 0x33)
SURFACE_2 = RGBColor(0x1F, 0x2A, 0x44)
INK = RGBColor(0xE7, 0xEC, 0xF5)          # texto principal
DIM = RGBColor(0x9A, 0xA5, 0xBE)          # secundario
ACCENT = RGBColor(0x7C, 0x9C, 0xFF)       # azul vibrante (apoyo)
ACCENT_2 = RGBColor(0x57, 0xD0, 0xB2)     # verde menta
WARN = RGBColor(0xF0, 0xB4, 0x4B)         # ámbar
HOT = RGBColor(0xE8, 0x6B, 0x7C)          # rojo coral (acento agudo)

# 16x9 widescreen
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ---------- Helpers ----------
def add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=18,
    bold=False,
    italic=False,
    color=INK,
    font="Calibri",
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def rich_textbox(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs: list of paragraphs; each is list of (text, opts) where opts is dict."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for pi, paragraph in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for text, opts in paragraph:
            r = p.add_run()
            r.text = text
            r.font.name = opts.get("font", "Calibri")
            r.font.size = Pt(opts.get("size", 18))
            r.font.bold = opts.get("bold", False)
            r.font.italic = opts.get("italic", False)
            r.font.color.rgb = opts.get("color", INK)
    return tb


def rect(slide, x, y, w, h, *, fill=SURFACE, line=None, rounded=False, line_w=1.0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    if rounded:
        shp.adjustments[0] = 0.12
    # remove default text
    shp.text_frame.text = ""
    return shp


def chip(slide, x, y, text, *, bg=ACCENT, fg=BG, size=12, padding_x=Inches(0.18), padding_y=Inches(0.08)):
    # measure-ish: width based on text length
    w = Inches(max(0.8, 0.11 * len(text) + 0.4))
    h = Inches(0.32)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.5
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = padding_x
    tf.margin_right = padding_x
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    return shp, w


def line(slide, x1, y1, x2, y2, *, color=ACCENT, weight=1.5, dash=False):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    if dash:
        from pptx.oxml.ns import qn
        from lxml import etree
        spPr = ln.line._get_or_add_ln()
        prstDash = etree.SubElement(spPr, qn("a:prstDash"))
        prstDash.set("val", "dash")
    return ln


def arrow(slide, x1, y1, x2, y2, *, color=ACCENT, weight=1.5):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    # add arrowhead via XML
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = ln.line._get_or_add_ln()
    tail = etree.SubElement(spPr, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return ln


def header(slide, kicker, title, *, kicker_color=ACCENT):
    """Cabecera consistente: kicker + título grande + barra inferior."""
    textbox(
        slide,
        Inches(0.55),
        Inches(0.45),
        Inches(11),
        Inches(0.4),
        kicker.upper(),
        size=12,
        bold=True,
        color=kicker_color,
        font="Calibri",
    )
    textbox(
        slide,
        Inches(0.55),
        Inches(0.78),
        Inches(12.2),
        Inches(0.9),
        title,
        size=34,
        bold=True,
        color=INK,
        font="Georgia",
    )
    rect(slide, Inches(0.55), Inches(1.7), Inches(0.7), Emu(38100), fill=ACCENT)


def footer(slide, n, total):
    textbox(
        slide,
        Inches(0.55),
        Inches(7.05),
        Inches(8),
        Inches(0.3),
        "MemGPT · LLMs como sistemas operativos",
        size=10,
        color=DIM,
    )
    textbox(
        slide,
        Inches(11.8),
        Inches(7.05),
        Inches(1.0),
        Inches(0.3),
        f"{n} / {total}",
        size=10,
        color=DIM,
        align=PP_ALIGN.RIGHT,
    )


def dot(slide, x, y, *, color=ACCENT, size=Inches(0.18)):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


# ---------- Slide builders ----------
def slide_title(prs):
    s = add_blank(prs)
    fill_bg(s, BG)
    # Decoración: bandas diagonales
    band = rect(s, Inches(8.5), Inches(-1), Inches(7), Inches(10), fill=SURFACE)
    band.rotation = 12
    band2 = rect(s, Inches(9.5), Inches(0.5), Inches(6), Inches(8), fill=SURFACE_2)
    band2.rotation = 12

    # Círculos decorativos
    o1 = slide_oval(s, Inches(10.8), Inches(1.3), Inches(2.6), ACCENT)
    o2 = slide_oval(s, Inches(9.2), Inches(4.8), Inches(1.6), ACCENT_2)
    o3 = slide_oval(s, Inches(11.6), Inches(5.2), Inches(0.9), HOT)

    # Kicker
    textbox(s, Inches(0.7), Inches(1.3), Inches(8), Inches(0.5),
            "CHARLA TÉCNICA · 2026", size=14, bold=True, color=ACCENT, font="Calibri")
    # Title
    textbox(s, Inches(0.7), Inches(1.8), Inches(9), Inches(2.5),
            "MemGPT", size=96, bold=True, color=INK, font="Georgia")
    textbox(s, Inches(0.7), Inches(3.6), Inches(9), Inches(1.5),
            "LLMs como sistemas operativos", size=36, color=DIM, font="Georgia", line_spacing=1.1)

    # Subtítulo
    textbox(s, Inches(0.7), Inches(5.0), Inches(8), Inches(0.8),
            "Memoria jerárquica auto-gestionada por el agente",
            size=20, color=INK, font="Calibri")

    # Footer
    rect(s, Inches(0.7), Inches(6.4), Inches(0.5), Emu(38100), fill=ACCENT)
    textbox(s, Inches(0.7), Inches(6.55), Inches(8), Inches(0.4),
            "Implementación sobre LangGraph · Pavel Kim 2023 · arXiv:2310.08560",
            size=14, color=DIM)
    textbox(s, Inches(0.7), Inches(6.95), Inches(8), Inches(0.4),
            "Máximo Fernández · maximofn.com",
            size=14, color=ACCENT, bold=True)


def slide_oval(slide, x, y, d, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def slide_problem(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "El problema", "La ventana de contexto se llena")

    # Dos columnas
    rect(s, Inches(0.55), Inches(2.1), Inches(6), Inches(4.6), fill=SURFACE, rounded=True)
    textbox(s, Inches(0.85), Inches(2.3), Inches(5.4), Inches(0.6),
            "Síntoma", size=18, bold=True, color=ACCENT_2)
    textbox(s, Inches(0.85), Inches(2.85), Inches(5.4), Inches(3.8),
            "Un LLM tiene contexto finito.\n\n"
            "Conversaciones largas, documentos grandes o tareas\n"
            "de varios días desbordan inevitablemente.\n\n"
            "Cuando se trunca, el agente olvida\n"
            "decisiones, identidad y hechos clave.",
            size=18, color=INK, line_spacing=1.35)

    rect(s, Inches(6.8), Inches(2.1), Inches(6), Inches(4.6), fill=SURFACE_2, rounded=True)
    textbox(s, Inches(7.1), Inches(2.3), Inches(5.4), Inches(0.6),
            "La idea de Karpathy (2023)", size=18, bold=True, color=ACCENT)
    textbox(s, Inches(7.1), Inches(2.85), Inches(5.4), Inches(3.8),
            "“El LLM es la CPU; la ventana de contexto es la RAM.”\n\n"
            "Y si es una CPU… necesita un sistema operativo:\n"
            "memoria virtual, planificador, llamadas al sistema.\n\n"
            "MemGPT lo construye.",
            size=18, color=INK, line_spacing=1.35, font="Georgia")
    footer(s, n, total)


def slide_os_analogy(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Analogía", "Memoria virtual para LLMs")

    # Dos bloques: SO vs MemGPT
    cols = [
        ("Sistema operativo", [
            ("CPU", "Ejecuta instrucciones"),
            ("RAM", "Pequeña, rápida"),
            ("Disco", "Grande, lento"),
            ("Paginación", "Mueve datos RAM ↔ disco"),
            ("Syscalls", "Procesos piden al SO"),
        ], ACCENT),
        ("MemGPT", [
            ("LLM", "Ejecuta turnos"),
            ("Ventana de contexto", "Lo que ve en cada inferencia"),
            ("Almacenamiento externo", "Recall / Archival / MemFS"),
            ("Queue Manager", "Expulsa y resume mensajes"),
            ("Function calls", "El LLM pide leer/escribir memoria"),
        ], ACCENT_2),
    ]
    x0 = Inches(0.55)
    col_w = Inches(6)
    for ci, (title, rows, c) in enumerate(cols):
        cx = x0 + (col_w + Inches(0.2)) * ci
        rect(s, cx, Inches(2.1), col_w, Inches(4.7), fill=SURFACE, rounded=True)
        textbox(s, cx + Inches(0.35), Inches(2.25), col_w - Inches(0.4), Inches(0.5),
                title, size=20, bold=True, color=c, font="Georgia")
        for ri, (k, v) in enumerate(rows):
            ry = Inches(2.9 + ri * 0.75)
            dot(s, cx + Inches(0.35), ry + Inches(0.08), color=c, size=Inches(0.14))
            textbox(s, cx + Inches(0.65), ry, Inches(2.0), Inches(0.5),
                    k, size=15, bold=True, color=INK)
            textbox(s, cx + Inches(2.7), ry, col_w - Inches(3.0), Inches(0.5),
                    v, size=14, color=DIM)
    footer(s, n, total)


def slide_three_pillars(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Qué es MemGPT", "Tres piezas, un agente que se administra solo")

    items = [
        ("01", "Jerarquía de memoria",
         "Rápida y pequeña (en contexto) +\nlenta e ilimitada (fuera).", ACCENT),
        ("02", "Gestor de control de flujo",
         "Decide cuándo invocar al LLM y\ncuándo expulsar y resumir mensajes.", ACCENT_2),
        ("03", "Funciones expuestas al LLM",
         "El propio modelo lee, escribe\ny busca en su memoria.", WARN),
    ]
    x0 = Inches(0.55)
    card_w = Inches(4.05)
    for i, (num, title, body, c) in enumerate(items):
        x = x0 + (card_w + Inches(0.15)) * i
        rect(s, x, Inches(2.2), card_w, Inches(4.3), fill=SURFACE, rounded=True)
        # acento lateral
        rect(s, x, Inches(2.2), Inches(0.18), Inches(4.3), fill=c)
        textbox(s, x + Inches(0.45), Inches(2.4), card_w - Inches(0.6), Inches(0.6),
                num, size=42, bold=True, color=c, font="Georgia")
        textbox(s, x + Inches(0.45), Inches(3.3), card_w - Inches(0.6), Inches(0.9),
                title, size=22, bold=True, color=INK, font="Georgia", line_spacing=1.1)
        textbox(s, x + Inches(0.45), Inches(4.4), card_w - Inches(0.6), Inches(2.0),
                body, size=15, color=DIM, line_spacing=1.4)

    rect(s, Inches(0.55), Inches(6.7), Inches(12.3), Inches(0.3), fill=SURFACE_2, rounded=True)
    textbox(s, Inches(0.85), Inches(6.7), Inches(11.6), Inches(0.3),
            "Punto clave: el LLM es auto-dirigido — decide qué guardar y qué buscar.",
            size=12, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, n, total)


def slide_context_window(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Anatomía", "Cómo modifica MemGPT la ventana de contexto")

    # Ventana grande con 3 bloques apilados
    win_x, win_y, win_w = Inches(1.0), Inches(2.0), Inches(7.3)
    rect(s, win_x - Inches(0.15), win_y - Inches(0.15),
         win_w + Inches(0.3), Inches(4.8), fill=SURFACE_2, rounded=True)

    blocks = [
        ("System Instructions", "Prompt + memGPT + funciones · solo lectura", ACCENT, Inches(0.9)),
        ("Working Context", "Bloques etiquetados editables (assistant / human / custom)", ACCENT_2, Inches(1.0)),
        ("FIFO Queue", "Slot 0: resumen recursivo  •  mensajes en orden cronológico", WARN, Inches(2.4)),
    ]
    cy = win_y
    for title, body, c, h in blocks:
        rect(s, win_x, cy, win_w, h, fill=SURFACE, rounded=True)
        rect(s, win_x, cy, Inches(0.16), h, fill=c)
        textbox(s, win_x + Inches(0.35), cy + Inches(0.15), win_w - Inches(0.5), Inches(0.4),
                title, size=17, bold=True, color=c, font="Georgia")
        textbox(s, win_x + Inches(0.35), cy + Inches(0.6), win_w - Inches(0.5), h - Inches(0.6),
                body, size=14, color=INK, line_spacing=1.35)
        cy += h + Inches(0.1)

    # Etiqueta "ventana de contexto"
    textbox(s, win_x, Inches(6.85), win_w, Inches(0.4),
            "Ventana de contexto del LLM", size=12, bold=True, color=DIM, align=PP_ALIGN.CENTER)

    # Columna derecha: leyenda + comparación
    rx = Inches(8.7)
    rw = Inches(4.2)
    rect(s, rx, Inches(2.0), rw, Inches(4.8), fill=SURFACE, rounded=True)
    textbox(s, rx + Inches(0.3), Inches(2.15), rw - Inches(0.4), Inches(0.45),
            "Agente normal", size=16, bold=True, color=DIM, font="Georgia")
    textbox(s, rx + Inches(0.3), Inches(2.55), rw - Inches(0.4), Inches(1.4),
            "Window:\n  [system prompt]\n  [user/agent turns…]",
            size=12, color=DIM, font="Consolas", line_spacing=1.25)

    textbox(s, rx + Inches(0.3), Inches(4.05), rw - Inches(0.4), Inches(0.45),
            "MemGPT", size=16, bold=True, color=ACCENT, font="Georgia")
    textbox(s, rx + Inches(0.3), Inches(4.45), rw - Inches(0.4), Inches(2.3),
            "Window:\n  [system + memGPT prompt]\n  [working context]\n  [recursive summary]\n  [user/agent + tool calls]",
            size=12, color=INK, font="Consolas", line_spacing=1.25)
    footer(s, n, total)


def slide_working_context(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Working Context", "La libreta del agente, siempre visible")

    # Izquierda: definición + propiedades
    textbox(s, Inches(0.55), Inches(2.05), Inches(6.5), Inches(0.5),
            "Bloque pequeño en el prompt, editable solo por el LLM vía tools.",
            size=18, color=INK, bold=False, font="Calibri")

    props = [
        ("Tamaño fijo por bloque", "Presupuesto de tokens; no crece."),
        ("Read / Write", "El LLM lo modifica; el harness, jamás."),
        ("Solo vía funciones", "core_memory_append · core_memory_replace"),
        ("Texto libre", "Lenguaje natural, no JSON."),
        ("Persistente", "Sobrevive entre sesiones."),
    ]
    for i, (k, v) in enumerate(props):
        py = Inches(2.85 + i * 0.62)
        dot(s, Inches(0.55), py + Inches(0.1), color=ACCENT_2, size=Inches(0.13))
        textbox(s, Inches(0.85), py, Inches(2.6), Inches(0.5),
                k, size=14, bold=True, color=INK)
        textbox(s, Inches(3.45), py, Inches(3.7), Inches(0.5),
                v, size=14, color=DIM)

    # Derecha: tarjetas de bloques etiquetados
    title = textbox(s, Inches(7.5), Inches(2.05), Inches(5.3), Inches(0.5),
            "Bloques típicos", size=18, bold=True, color=ACCENT, font="Georgia")
    blocks = [
        ("assistant", "Identidad del agente: nombre, personalidad, rol.",
         "“Soy Sam, paciente y curiosa. No revelo que soy IA.”", ACCENT),
        ("human", "Lo que el agente recuerda del usuario.",
         "“Ana, vive en Madrid, le gusta el surf.”", ACCENT_2),
        ("custom", "Bloques que el desarrollador añade.",
         "task_context, project_constraints, team_members…", WARN),
    ]
    for i, (label, desc, ex, c) in enumerate(blocks):
        y = Inches(2.65 + i * 1.45)
        rect(s, Inches(7.5), y, Inches(5.3), Inches(1.3), fill=SURFACE, rounded=True)
        chip(s, Inches(7.7), y + Inches(0.15), label, bg=c, size=11)
        textbox(s, Inches(7.7), y + Inches(0.55), Inches(4.9), Inches(0.4),
                desc, size=13, color=INK)
        textbox(s, Inches(7.7), y + Inches(0.9), Inches(4.9), Inches(0.4),
                ex, size=12, color=DIM, font="Consolas")
    footer(s, n, total)


def slide_core_memory_tools(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Tools de Core Memory", "Cómo edita el LLM su libreta")

    examples = [
        ("core_memory_append",
         "Añade información al final del bloque.",
         'core_memory_append(\n  label="human",\n  content="Cumpleaños: 7 de febrero"\n)', ACCENT_2),
        ("core_memory_replace",
         "Sustituye un fragmento (corrige, actualiza).",
         'core_memory_replace(\n  label="human",\n  old="novio James",\n  new="ex-novio James"\n)', WARN),
    ]
    for i, (name, desc, code, c) in enumerate(examples):
        x = Inches(0.55) + Inches(6.25) * i
        rect(s, x, Inches(2.1), Inches(6.0), Inches(4.4), fill=SURFACE, rounded=True)
        rect(s, x, Inches(2.1), Inches(6.0), Inches(0.6), fill=SURFACE_2)
        textbox(s, x + Inches(0.3), Inches(2.18), Inches(5.5), Inches(0.5),
                name, size=18, bold=True, color=c, font="Consolas")
        textbox(s, x + Inches(0.3), Inches(2.85), Inches(5.5), Inches(0.6),
                desc, size=15, color=INK)
        rect(s, x + Inches(0.3), Inches(3.6), Inches(5.4), Inches(2.7), fill=BG, rounded=True)
        textbox(s, x + Inches(0.55), Inches(3.78), Inches(5.0), Inches(2.5),
                code, size=15, color=INK, font="Consolas", line_spacing=1.35)

    rect(s, Inches(0.55), Inches(6.7), Inches(12.3), Inches(0.3), fill=SURFACE_2, rounded=True)
    textbox(s, Inches(0.85), Inches(6.7), Inches(11.6), Inches(0.3),
            "El LLM nunca escribe el bloque generando texto: solo a través de estas funciones.",
            size=12, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, n, total)


def slide_fifo(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "FIFO Queue", "El flujo en vivo de la conversación")

    # Diagrama de cola: bloques apilados
    qx, qy, qw = Inches(0.6), Inches(2.1), Inches(7.2)
    rect(s, qx - Inches(0.1), qy - Inches(0.1), qw + Inches(0.2), Inches(4.7),
         fill=SURFACE_2, rounded=True)

    entries = [
        ("0", "Resumen recursivo", "“Lo que pasó antes”, regenerado al hacer flush.", HOT),
        ("1", "(user)", "Hola, ¿te acuerdas dónde comimos?", INK),
        ("2", "(assistant)", "tool_call: recall_memory_search(query='comida pacifica')", ACCENT),
        ("3", "(tool_result)", "[03/12] “Taco Bell junto a la playa”", ACCENT_2),
        ("4", "(assistant)", "Sí, en el Taco Bell de Pacifica.", INK),
        ("5", "(system)", "Memory Pressure: contexto al 75%", WARN),
        ("…", "(expulsados)", "Se mueven a Recall Storage", DIM),
    ]
    for i, (idx, role, txt, c) in enumerate(entries):
        y = qy + Inches(0.1 + i * 0.62)
        rect(s, qx + Inches(0.1), y, qw - Inches(0.2), Inches(0.55),
             fill=SURFACE if i % 2 == 0 else BG, rounded=False)
        textbox(s, qx + Inches(0.2), y + Inches(0.13), Inches(0.4), Inches(0.4),
                idx, size=11, bold=True, color=DIM, align=PP_ALIGN.CENTER)
        textbox(s, qx + Inches(0.65), y + Inches(0.13), Inches(1.6), Inches(0.4),
                role, size=12, bold=True, color=c, font="Consolas")
        textbox(s, qx + Inches(2.3), y + Inches(0.13), qw - Inches(2.6), Inches(0.4),
                txt, size=12, color=INK)

    # Etiqueta extremos
    textbox(s, qx, Inches(6.7), Inches(3.5), Inches(0.25),
            "↑ más antiguos (se expulsan)", size=10, color=DIM)
    textbox(s, qx + qw - Inches(3.5), Inches(6.7), Inches(3.5), Inches(0.25),
            "más recientes ↓", size=10, color=DIM, align=PP_ALIGN.RIGHT)

    # Columna derecha: qué se guarda
    rx = Inches(8.2)
    rect(s, rx, Inches(2.0), Inches(4.7), Inches(4.8), fill=SURFACE, rounded=True)
    textbox(s, rx + Inches(0.3), Inches(2.15), Inches(4.2), Inches(0.5),
            "Entra en la FIFO", size=18, bold=True, color=ACCENT, font="Georgia")
    items = [
        "Mensajes de usuario y agente",
        "Tool calls de MemGPT y otras tools",
        "Resultados de las funciones",
        "Mensajes del sistema (alertas, eventos)",
        "Sellos temporales para anclar al agente",
    ]
    for i, t in enumerate(items):
        y = Inches(2.75 + i * 0.55)
        dot(s, rx + Inches(0.35), y + Inches(0.1), color=ACCENT, size=Inches(0.12))
        textbox(s, rx + Inches(0.65), y, Inches(3.9), Inches(0.5),
                t, size=14, color=INK)
    textbox(s, rx + Inches(0.3), Inches(5.85), Inches(4.2), Inches(0.9),
            "Lo escribe el Queue Manager, no el LLM:\nes un log automático.",
            size=12, color=DIM, line_spacing=1.3, italic=False)
    footer(s, n, total)


def slide_thresholds(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Umbrales", "70 % avisa, 100 % expulsa")

    # Barra horizontal del contexto
    bar_x, bar_y, bar_w, bar_h = Inches(1.0), Inches(3.0), Inches(11.3), Inches(0.9)
    rect(s, bar_x, bar_y, bar_w, bar_h, fill=SURFACE_2, rounded=True)
    # rellenos por tramos
    green = bar_w * 0.7
    amber = bar_w * 0.3
    rect(s, bar_x, bar_y, int(green), bar_h, fill=ACCENT_2, rounded=True)
    rect(s, bar_x + int(green), bar_y, int(amber), bar_h, fill=WARN, rounded=False)

    # Marcadores
    def marker(x, label, sub, color, sub_offset=Inches(1.5), sub_align=PP_ALIGN.CENTER):
        line(s, x, bar_y - Inches(0.2), x, bar_y + bar_h + Inches(0.2),
             color=color, weight=2)
        textbox(s, x - Inches(0.8), bar_y - Inches(0.95), Inches(1.6), Inches(0.4),
                label, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        textbox(s, x - sub_offset, bar_y + bar_h + Inches(0.25), Inches(3.0), Inches(0.5),
                sub, size=12, color=DIM, align=sub_align)

    marker(bar_x + int(green), "70 %", "Warning · Memory Pressure Alert", WARN)
    marker(bar_x + bar_w, "100 %", "Flush · expulsa y resume",
           HOT, sub_offset=Inches(2.7), sub_align=PP_ALIGN.RIGHT)

    # Cajas inferiores
    rect(s, Inches(0.55), Inches(5.1), Inches(6.0), Inches(1.8), fill=SURFACE, rounded=True)
    textbox(s, Inches(0.85), Inches(5.25), Inches(5.6), Inches(0.5),
            "Warning (70 %)", size=18, bold=True, color=WARN, font="Georgia")
    textbox(s, Inches(0.85), Inches(5.7), Inches(5.6), Inches(1.1),
            "Se inyecta una alerta en la FIFO.\nEl LLM consolida lo importante en\nWorking Context, Archival o MemFS.",
            size=14, color=INK, line_spacing=1.35)

    rect(s, Inches(6.85), Inches(5.1), Inches(6.0), Inches(1.8), fill=SURFACE, rounded=True)
    textbox(s, Inches(7.15), Inches(5.25), Inches(5.6), Inches(0.5),
            "Flush (100 %)", size=18, bold=True, color=HOT, font="Georgia")
    textbox(s, Inches(7.15), Inches(5.7), Inches(5.6), Inches(1.1),
            "El Queue Manager expulsa los mensajes\nmás viejos (~50 %) y regenera el resumen.\nNada se pierde: vive en Recall Storage.",
            size=14, color=INK, line_spacing=1.35)
    footer(s, n, total)


def slide_recursive_summary(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Resumen recursivo", "Una llamada al LLM fuera del bucle principal")

    # Diagrama: 3 inputs → summarizer → 2 outputs
    in_x, in_w = Inches(0.55), Inches(3.6)
    inputs = [
        ("Working Context actual", ACCENT_2),
        ("Resumen previo (slot 0)", WARN),
        ("Mensajes a expulsar", HOT),
    ]
    for i, (t, c) in enumerate(inputs):
        y = Inches(2.5 + i * 1.15)
        rect(s, in_x, y, in_w, Inches(0.9), fill=SURFACE, rounded=True)
        rect(s, in_x, y, Inches(0.16), Inches(0.9), fill=c)
        textbox(s, in_x + Inches(0.35), y + Inches(0.25), in_w - Inches(0.5), Inches(0.5),
                t, size=15, bold=True, color=INK)
        arrow(s, in_x + in_w, y + Inches(0.45),
              Inches(5.55), Inches(4.05), color=ACCENT, weight=1.5)

    # Summarizer central
    rect(s, Inches(5.55), Inches(3.4), Inches(2.4), Inches(1.3),
         fill=ACCENT, rounded=True)
    textbox(s, Inches(5.55), Inches(3.5), Inches(2.4), Inches(0.5),
            "Summarizer LLM", size=16, bold=True, color=BG, align=PP_ALIGN.CENTER, font="Georgia")
    textbox(s, Inches(5.55), Inches(3.95), Inches(2.4), Inches(0.8),
            "Llamada independiente,\nprompt mínimo",
            size=12, color=BG, align=PP_ALIGN.CENTER, line_spacing=1.3)

    # Outputs
    outs = [("Nuevo resumen recursivo", ACCENT),
            ("Promociones a Working Context", ACCENT_2)]
    for i, (t, c) in enumerate(outs):
        y = Inches(3.1 + i * 1.15)
        rect(s, Inches(8.6), y, Inches(4.3), Inches(0.9), fill=SURFACE, rounded=True)
        rect(s, Inches(8.6) + Inches(4.3) - Inches(0.16), y, Inches(0.16), Inches(0.9), fill=c)
        textbox(s, Inches(8.85), y + Inches(0.25), Inches(3.7), Inches(0.5),
                t, size=15, bold=True, color=INK)
        arrow(s, Inches(7.95), Inches(4.05), Inches(8.6), y + Inches(0.45),
              color=ACCENT, weight=1.5)

    # Nota inferior
    rect(s, Inches(0.55), Inches(6.05), Inches(12.3), Inches(0.85),
         fill=SURFACE_2, rounded=True)
    textbox(s, Inches(0.85), Inches(6.15), Inches(11.8), Inches(0.7),
            "¿Por qué fuera del bucle? Si el contexto está al 100 %, no cabe ‘hazme un resumen’ encima.\n"
            "El summarizer recibe su propio prompt, sin las System Instructions del agente principal.",
            size=13, color=INK, line_spacing=1.3)
    footer(s, n, total)


def slide_recall(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Recall Storage", "El log completo de la conversación")

    rect(s, Inches(0.55), Inches(2.1), Inches(6.3), Inches(4.6), fill=SURFACE, rounded=True)
    textbox(s, Inches(0.85), Inches(2.3), Inches(5.7), Inches(0.5),
            "¿Qué es?", size=18, bold=True, color=ACCENT, font="Georgia")
    textbox(s, Inches(0.85), Inches(2.8), Inches(5.7), Inches(3.7),
            "Historial completo, fuera del contexto.\n\n"
            "Cada mensaje que entra a la FIFO se copia\n"
            "automáticamente a Recall en el mismo turno.\n\n"
            "Cuando se hace flush, nada se pierde:\n"
            "los mensajes expulsados ya están en Recall.",
            size=15, color=INK, line_spacing=1.4)

    rect(s, Inches(7.15), Inches(2.1), Inches(5.7), Inches(4.6), fill=SURFACE_2, rounded=True)
    textbox(s, Inches(7.45), Inches(2.3), Inches(5.1), Inches(0.5),
            "Cómo lo usa el LLM", size=18, bold=True, color=ACCENT_2, font="Georgia")
    rect(s, Inches(7.45), Inches(2.85), Inches(5.1), Inches(0.75), fill=BG, rounded=True)
    textbox(s, Inches(7.6), Inches(2.95), Inches(5.0), Inches(0.6),
            'recall_memory_search(query="...", page=N)',
            size=14, color=INK, font="Consolas")
    props = [
        ("Solo lectura", "el LLM busca, no inserta"),
        ("Búsqueda semántica", "embeddings + similitud"),
        ("Paginada", "evita ‘lost in the middle’"),
        ("Ilimitada", "toda la historia"),
        ("Persistente", "entre sesiones"),
    ]
    for i, (k, v) in enumerate(props):
        y = Inches(3.85 + i * 0.55)
        dot(s, Inches(7.45), y + Inches(0.1), color=ACCENT_2, size=Inches(0.12))
        textbox(s, Inches(7.75), y, Inches(2.0), Inches(0.5),
                k, size=13, bold=True, color=INK)
        textbox(s, Inches(9.7), y, Inches(3.1), Inches(0.5),
                v, size=13, color=DIM)
    footer(s, n, total)


def slide_archival(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Archival Storage", "Lo que el agente decide guardar a largo plazo")

    # tres columnas: qué, escritura, lectura
    cards = [
        ("¿Qué entra?",
         "Texto arbitrario que el LLM elige persistir:\n"
         "• resúmenes propios\n"
         "• hechos sueltos\n"
         "• documentos del usuario\n"
         "• conclusiones de la sesión",
         ACCENT),
        ("Escritura explícita",
         "archival_memory_insert(\n  content=\"…\"\n)\n\n"
         "Solo entra lo que el LLM\ndecide insertar. Recall, no.",
         ACCENT_2),
        ("Lectura iterativa",
         "archival_memory_search(\n  query=\"…\", page=N\n)\n\n"
         "Embeddings + paginación.\nEl LLM repite, reformula, pagina.",
         WARN),
    ]
    x0 = Inches(0.55)
    cw = Inches(4.05)
    for i, (title, body, c) in enumerate(cards):
        x = x0 + (cw + Inches(0.15)) * i
        rect(s, x, Inches(2.1), cw, Inches(4.7), fill=SURFACE, rounded=True)
        rect(s, x, Inches(2.1), cw, Inches(0.6), fill=c)
        textbox(s, x + Inches(0.3), Inches(2.18), cw - Inches(0.4), Inches(0.5),
                title, size=17, bold=True, color=BG, font="Georgia")
        is_code = "(" in body and "=" in body
        textbox(s, x + Inches(0.3), Inches(2.95), cw - Inches(0.4), Inches(3.7),
                body, size=14 if not is_code else 13, color=INK,
                font="Consolas" if is_code else "Calibri",
                line_spacing=1.4)
    footer(s, n, total)


def slide_archival_vs_rag(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Archival vs RAG clásico", "Mismo motor, distinto director")

    rows = [
        ("Quién decide buscar", "El pipeline, siempre", "El LLM, solo si lo necesita"),
        ("Iteración", "Un disparo", "Paginar, reformular, reintentar"),
        ("Escritura", "Corpus estático",
         "El LLM puede archival_memory_insert"),
        ("Resultado top-K", "Inyectado al prompt",
         "El LLM lee y decide qué hacer"),
    ]
    # encabezados
    x0 = Inches(0.55)
    col_w = [Inches(3.3), Inches(4.6), Inches(4.6)]
    cx = x0
    headers = [("", BG), ("RAG clásico", DIM), ("Archival Storage", ACCENT)]
    for i, (h, c) in enumerate(headers):
        rect(s, cx, Inches(2.1), col_w[i], Inches(0.6),
             fill=SURFACE_2 if i else BG, rounded=False)
        textbox(s, cx + Inches(0.25), Inches(2.2), col_w[i] - Inches(0.3), Inches(0.5),
                h, size=15, bold=True, color=c if i else INK, font="Georgia")
        cx += col_w[i]

    for ri, (k, a, b) in enumerate(rows):
        ry = Inches(2.78 + ri * 0.85)
        cx = x0
        rect(s, cx, ry, sum(col_w, Inches(0)), Inches(0.8),
             fill=SURFACE if ri % 2 == 0 else BG, rounded=False)
        textbox(s, cx + Inches(0.25), ry + Inches(0.2), col_w[0] - Inches(0.3), Inches(0.5),
                k, size=14, bold=True, color=INK)
        cx += col_w[0]
        textbox(s, cx + Inches(0.25), ry + Inches(0.2), col_w[1] - Inches(0.3), Inches(0.5),
                a, size=13, color=DIM)
        cx += col_w[1]
        textbox(s, cx + Inches(0.25), ry + Inches(0.2), col_w[2] - Inches(0.3), Inches(0.5),
                b, size=13, color=INK)

    rect(s, Inches(0.55), Inches(6.3), Inches(12.3), Inches(0.7),
         fill=SURFACE_2, rounded=True)
    textbox(s, Inches(0.85), Inches(6.42), Inches(11.7), Inches(0.5),
            "Archival = RAG (motor) + control del LLM (cuándo, cómo, cuántas veces) + escritura por el agente.",
            size=14, bold=True, color=ACCENT)
    footer(s, n, total)


def slide_memfs(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "MemFS", "Memoria como filesystem versionado")

    # Izquierda: árbol
    rect(s, Inches(0.55), Inches(2.1), Inches(5.5), Inches(4.7),
         fill=SURFACE, rounded=True)
    textbox(s, Inches(0.85), Inches(2.25), Inches(5.0), Inches(0.5),
            "/", size=14, bold=True, color=ACCENT, font="Consolas")
    tree = [
        ("├── users/", DIM),
        ("│   └── maximo/", DIM),
        ("│       └── objetivos-q2.md", INK),
        ("├── projects/", DIM),
        ("│   └── portafolio/", DIM),
        ("│       ├── decisiones.md", INK),
        ("│       └── notas.md", INK),
        ("└── learning/", DIM),
        ("    └── spanish/", DIM),
    ]
    for i, (t, c) in enumerate(tree):
        textbox(s, Inches(0.85), Inches(2.7 + i * 0.42),
                Inches(5.0), Inches(0.4),
                t, size=14, color=c, font="Consolas")

    # Derecha: tools y ventajas
    rect(s, Inches(6.3), Inches(2.1), Inches(6.55), Inches(2.6),
         fill=SURFACE_2, rounded=True)
    textbox(s, Inches(6.55), Inches(2.25), Inches(6.0), Inches(0.5),
            "Tools de MemFS", size=17, bold=True, color=ACCENT, font="Georgia")
    code = ("memfs.create(path, content)\n"
            "memfs.read(path)        memfs.write(path, content)\n"
            "memfs.list(path)        memfs.move(src, dst)\n"
            "memfs.grep(path, q)     memfs.history(path)\n"
            "memfs.rollback(path, version)")
    textbox(s, Inches(6.55), Inches(2.8), Inches(6.0), Inches(2.0),
            code, size=13, color=INK, font="Consolas", line_spacing=1.4)

    rect(s, Inches(6.3), Inches(4.85), Inches(6.55), Inches(1.95),
         fill=SURFACE, rounded=True)
    textbox(s, Inches(6.55), Inches(4.95), Inches(6.0), Inches(0.5),
            "Por qué versionado", size=17, bold=True, color=ACCENT_2, font="Georgia")
    advs = [
        "Proyectos largos en árbol jerárquico",
        "Trazabilidad: cómo evolucionó un fichero",
        "Rollback si el agente sobrescribe por error",
        "Pizarra compartida entre multi-agentes",
    ]
    for i, t in enumerate(advs):
        y = Inches(5.48 + i * 0.32)
        dot(s, Inches(6.6), y + Inches(0.05), color=ACCENT_2, size=Inches(0.11))
        textbox(s, Inches(6.85), y, Inches(5.7), Inches(0.4),
                t, size=12, color=INK)
    footer(s, n, total)


def slide_compare_layers(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Working / Archival / MemFS", "Tres memorias que coexisten")

    headers_ = [("", BG, Inches(2.6)),
                ("Working Context", ACCENT, Inches(3.3)),
                ("Archival Storage", ACCENT_2, Inches(3.5)),
                ("MemFS", WARN, Inches(3.0))]
    rows = [
        ("Dónde vive", "En la ventana de contexto", "BD vectorial", "Filesystem versionado"),
        ("Tamaño", "Pequeño (tokens)", "Ilimitado", "Ilimitado"),
        ("Siempre visible", "Sí", "No (buscar)", "No (navegar)"),
        ("Búsqueda", "Directa", "Semántica", "Path + grep"),
        ("Versionado", "No", "No", "Sí, tipo git"),
        ("Coste por uso", "Tokens en cada turno", "Tool call", "Tool call"),
    ]

    # header row
    cx = Inches(0.55)
    y0 = Inches(2.05)
    for i, (h, c, w) in enumerate(headers_):
        rect(s, cx, y0, w, Inches(0.55),
             fill=BG if i == 0 else SURFACE_2, rounded=False)
        if i > 0:
            rect(s, cx, y0, w, Inches(0.08), fill=c)
        textbox(s, cx + Inches(0.2), y0 + Inches(0.12), w - Inches(0.25), Inches(0.4),
                h, size=14, bold=True, color=c if i else INK, font="Georgia",
                align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.LEFT)
        cx += w

    for ri, row in enumerate(rows):
        ry = Inches(2.7 + ri * 0.65)
        cx = Inches(0.55)
        rect(s, cx, ry, sum((w for _, _, w in headers_), Inches(0)),
             Inches(0.6), fill=SURFACE if ri % 2 == 0 else BG)
        for i, val in enumerate(row):
            w = headers_[i][2]
            textbox(s, cx + Inches(0.2), ry + Inches(0.15), w - Inches(0.25), Inches(0.4),
                    val, size=13,
                    bold=(i == 0),
                    color=INK if i == 0 else DIM)
            cx += w
    footer(s, n, total)


def slide_decision_matrix(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Decidir dónde guardar", "Regla del agente")

    cards = [
        ("Working Context",
         "Crítico, debe verse en cada turno.\n"
         "Identidad, preferencias clave,\nestado actual del usuario.",
         "“El usuario es Máximo,\nprefiere español, usa\nLangGraph + Astro.”",
         ACCENT),
        ("Archival Storage",
         "Voluminoso, no sabes cuándo lo necesitarás.\n"
         "Búsqueda por significado a posteriori.",
         "“El usuario eligió pgvector\nporque…” →\nrecuperable meses después.",
         ACCENT_2),
        ("MemFS",
         "Estructurado, evoluciona en el tiempo,\n"
         "necesitas historial.",
         "/projects/portafolio/\n  decisiones.md\n+ memfs.history(...)",
         WARN),
    ]
    x0 = Inches(0.55)
    cw = Inches(4.05)
    for i, (title, when, example, c) in enumerate(cards):
        x = x0 + (cw + Inches(0.15)) * i
        rect(s, x, Inches(2.05), cw, Inches(4.85), fill=SURFACE, rounded=True)
        rect(s, x, Inches(2.05), cw, Inches(0.7), fill=c)
        textbox(s, x + Inches(0.3), Inches(2.2), cw - Inches(0.4), Inches(0.5),
                title, size=18, bold=True, color=BG, font="Georgia")
        textbox(s, x + Inches(0.3), Inches(3.0), cw - Inches(0.4), Inches(0.4),
                "Si el dato es…", size=13, bold=True, color=DIM)
        textbox(s, x + Inches(0.3), Inches(3.4), cw - Inches(0.4), Inches(1.5),
                when, size=14, color=INK, line_spacing=1.35)
        textbox(s, x + Inches(0.3), Inches(5.05), cw - Inches(0.4), Inches(0.4),
                "Ejemplo", size=13, bold=True, color=DIM)
        rect(s, x + Inches(0.3), Inches(5.45), cw - Inches(0.6), Inches(1.4),
             fill=BG, rounded=True)
        textbox(s, x + Inches(0.45), Inches(5.55), cw - Inches(0.9), Inches(1.25),
                example, size=12, color=INK, font="Consolas", line_spacing=1.3)
    footer(s, n, total)


def slide_cycle(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Ciclo de funcionamiento", "Qué pasa cuando llega un evento")

    # Pipeline horizontal con 5 cajas
    boxes = [
        ("Evento", "Mensaje, alerta o tick", ACCENT),
        ("Queue Manager", "Añade a FIFO + Recall", ACCENT_2),
        ("Pressure check", "¿70 % / 100 %?", WARN),
        ("Inferencia LLM", "Prompt = system + Core + FIFO", ACCENT),
        ("¿Tool call?", "Sí → ejecutar y repetir\nNo → respuesta", HOT),
    ]
    n_b = len(boxes)
    pad = Inches(0.25)
    total_w = SLIDE_W - Inches(1.1)
    bw = (total_w - pad * (n_b - 1)) / n_b
    y = Inches(2.6)
    bh = Inches(2.0)
    for i, (t, sub, c) in enumerate(boxes):
        x = Inches(0.55) + (bw + pad) * i
        rect(s, x, y, bw, bh, fill=SURFACE, rounded=True)
        rect(s, x, y, bw, Inches(0.4), fill=c)
        textbox(s, x + Inches(0.1), y + Inches(0.4), bw - Inches(0.2), Inches(0.6),
                f"{i+1:02d}", size=22, bold=True, color=c, font="Georgia",
                align=PP_ALIGN.CENTER)
        textbox(s, x + Inches(0.15), y + Inches(1.0), bw - Inches(0.3), Inches(0.5),
                t, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER, font="Georgia")
        textbox(s, x + Inches(0.15), y + Inches(1.4), bw - Inches(0.3), Inches(0.55),
                sub, size=11, color=DIM, align=PP_ALIGN.CENTER, line_spacing=1.25)
        if i < n_b - 1:
            arrow(s, x + bw, y + bh / 2, x + bw + pad, y + bh / 2,
                  color=ACCENT, weight=2.2)

    # Loop arrow back
    line(s, Inches(0.55) + bw * 0.5, y + bh + Inches(0.15),
         Inches(0.55) + bw * 0.5, y + bh + Inches(0.8),
         color=DIM, weight=1.2)
    line(s, Inches(0.55) + bw * 0.5, y + bh + Inches(0.8),
         Inches(0.55) + bw * (n_b - 0.5) + pad * (n_b - 1),
         y + bh + Inches(0.8), color=DIM, weight=1.2)
    arrow(s, Inches(0.55) + bw * (n_b - 0.5) + pad * (n_b - 1), y + bh + Inches(0.8),
          Inches(0.55) + bw * (n_b - 0.5) + pad * (n_b - 1), y + bh + Inches(0.15),
          color=DIM, weight=1.2)
    textbox(s, Inches(4.5), y + bh + Inches(0.95), Inches(4), Inches(0.5),
            "si hubo tool call → siguiente paso", size=12, color=DIM,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    footer(s, n, total)


def slide_events(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Eventos automáticos", "El agente puede arrancarse solo")

    cards = [
        ("Wall-clock",
         "El reloj del mundo dispara una acción.",
         "“Da los buenos días a las 09:00”\n"
         "“Manda el correo en 30 minutos”\n"
         "“Comprueba el estado cada 10 min”",
         ACCENT),
        ("Sleep-time agents",
         "Cada N pasos del agente principal\n"
         "se activa un agente auxiliar.",
         "“Cada 10 turnos consolida memoria”\n"
         "“Cada 20 mensajes revisa coherencia”\n"
         "“Tras 50 turnos sin búsqueda, sugiere\n recuperar de archival”",
         ACCENT_2),
    ]
    x0 = Inches(0.55)
    cw = Inches(6.15)
    for i, (title, desc, ex, c) in enumerate(cards):
        x = x0 + (cw + Inches(0.2)) * i
        rect(s, x, Inches(2.1), cw, Inches(4.7), fill=SURFACE, rounded=True)
        rect(s, x, Inches(2.1), cw, Inches(0.7), fill=c)
        textbox(s, x + Inches(0.3), Inches(2.25), cw - Inches(0.4), Inches(0.5),
                title, size=20, bold=True, color=BG, font="Georgia")
        textbox(s, x + Inches(0.3), Inches(3.0), cw - Inches(0.4), Inches(1.0),
                desc, size=15, color=INK, line_spacing=1.35)
        textbox(s, x + Inches(0.3), Inches(4.2), cw - Inches(0.4), Inches(0.4),
                "Ejemplos", size=13, bold=True, color=DIM)
        rect(s, x + Inches(0.3), Inches(4.65), cw - Inches(0.6), Inches(2.0),
             fill=BG, rounded=True)
        textbox(s, x + Inches(0.45), Inches(4.8), cw - Inches(0.9), Inches(1.8),
                ex, size=13, color=INK, font="Consolas", line_spacing=1.5)
    footer(s, n, total)


def slide_heartbeat(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Encadenar acciones", "request_heartbeat (LLMs antiguos) vs nativo")

    # Izquierda: legacy
    rect(s, Inches(0.55), Inches(2.1), Inches(6.0), Inches(4.7), fill=SURFACE, rounded=True)
    textbox(s, Inches(0.85), Inches(2.25), Inches(5.5), Inches(0.5),
            "LEGACY · request_heartbeat=True", size=16, bold=True, color=WARN, font="Georgia")
    textbox(s, Inches(0.85), Inches(2.85), Inches(5.5), Inches(2.3),
            "Argumento booleano en cualquier tool call.\n"
            "Si está activo, MemGPT vuelve a invocar al LLM\n"
            "aunque éste hubiera ‘parado’.\n\n"
            "Evita que el modelo se rinda a medio buscar.",
            size=14, color=INK, line_spacing=1.4)
    rect(s, Inches(0.85), Inches(5.0), Inches(5.4), Inches(1.7), fill=BG, rounded=True)
    textbox(s, Inches(1.05), Inches(5.15), Inches(5.0), Inches(1.5),
            'search_files(\n  query="MyClass",\n  request_heartbeat=True\n)',
            size=14, color=INK, font="Consolas", line_spacing=1.4)

    # Derecha: native + riesgo
    rect(s, Inches(6.85), Inches(2.1), Inches(6.0), Inches(2.2), fill=SURFACE, rounded=True)
    textbox(s, Inches(7.15), Inches(2.25), Inches(5.5), Inches(0.5),
            "NATIVE · LLMs modernos", size=16, bold=True, color=ACCENT_2, font="Georgia")
    textbox(s, Inches(7.15), Inches(2.85), Inches(5.5), Inches(1.4),
            "Si la respuesta contiene tool_calls, se ejecutan.\n"
            "Si no, es el mensaje final al usuario.\n"
            "No hace falta booleano: el LLM decide.",
            size=14, color=INK, line_spacing=1.4)

    rect(s, Inches(6.85), Inches(4.45), Inches(6.0), Inches(2.35), fill=SURFACE_2, rounded=True)
    textbox(s, Inches(7.15), Inches(4.6), Inches(5.5), Inches(0.5),
            "Protección contra bucles", size=16, bold=True, color=HOT, font="Georgia")
    items = [
        "Límite de heartbeats por turno",
        "Timeout wall-clock",
        "Detección de tool calls repetidas sin progreso",
    ]
    for i, t in enumerate(items):
        y = Inches(5.15 + i * 0.5)
        dot(s, Inches(7.15), y + Inches(0.1), color=HOT, size=Inches(0.13))
        textbox(s, Inches(7.45), y, Inches(5.2), Inches(0.5), t, size=13, color=INK)
    footer(s, n, total)


def slide_my_impl_overview(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Mi implementación", "MemGPT sobre LangGraph, desde cero", kicker_color=ACCENT_2)

    # Grafo simplificado
    boxes = [
        ("turn_init", Inches(0.6), Inches(2.3), ACCENT),
        ("recall_sync_in", Inches(2.55), Inches(2.3), ACCENT),
        ("pressure_check", Inches(4.7), Inches(2.3), WARN),
        ("agent (LLM)", Inches(6.9), Inches(2.3), ACCENT_2),
        ("step_tick", Inches(9.05), Inches(2.3), ACCENT),
        ("recall_sync_post", Inches(11.0), Inches(2.3), ACCENT),
    ]
    for (t, x, y, c) in boxes:
        rect(s, x, y, Inches(1.95), Inches(0.7), fill=SURFACE, rounded=True)
        rect(s, x, y, Inches(0.12), Inches(0.7), fill=c)
        textbox(s, x + Inches(0.2), y + Inches(0.18), Inches(1.7), Inches(0.4),
                t, size=12, bold=True, color=INK, font="Consolas", align=PP_ALIGN.LEFT)
    # arrows entre los principales
    for i in range(len(boxes) - 1):
        arrow(s, boxes[i][1] + Inches(1.95), boxes[i][2] + Inches(0.35),
              boxes[i + 1][1], boxes[i + 1][2] + Inches(0.35), color=ACCENT, weight=1.5)

    # tools loop
    rect(s, Inches(5.5), Inches(3.6), Inches(2.85), Inches(0.7), fill=SURFACE, rounded=True)
    rect(s, Inches(5.5), Inches(3.6), Inches(0.12), Inches(0.7), fill=HOT)
    textbox(s, Inches(5.7), Inches(3.78), Inches(2.6), Inches(0.4),
            "tools + heartbeat_check", size=12, bold=True, color=INK, font="Consolas")
    arrow(s, Inches(7.85), Inches(3.0), Inches(6.95), Inches(3.6),
          color=HOT, weight=1.5)
    arrow(s, Inches(5.5), Inches(3.85), Inches(3.5), Inches(3.0),
          color=HOT, weight=1.5)

    # Stack
    rect(s, Inches(0.55), Inches(4.7), Inches(12.3), Inches(2.15), fill=SURFACE_2, rounded=True)
    textbox(s, Inches(0.85), Inches(4.85), Inches(11.7), Inches(0.5),
            "Stack", size=18, bold=True, color=ACCENT, font="Georgia")
    stack = [
        ("LangGraph", "StateGraph + MemGPTState (Pydantic)", ACCENT),
        ("Postgres + pgvector", "PostgresSaver + EventStore", ACCENT_2),
        ("Neo4j + Graphiti", "Recall + Archival con embeddings", WARN),
        ("MemFS in-memory", "9 tools versionadas (snapshot + SHA-1)", HOT),
    ]
    for i, (k, v, c) in enumerate(stack):
        x = Inches(0.85) + Inches(3.0) * i
        chip(s, x, Inches(5.45), k, bg=c, size=11)
        textbox(s, x, Inches(5.95), Inches(2.85), Inches(0.85),
                v, size=12, color=INK, line_spacing=1.35)
    footer(s, n, total)


def slide_my_tools(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Tools del agente", "Lo que el LLM puede llamar")

    groups = [
        ("Core Memory", ACCENT, [
            "core_memory_append(label, content)",
            "core_memory_replace(label, old, new)",
        ]),
        ("Recall & Archival", ACCENT_2, [
            "recall_memory_search(query, page)",
            "archival_memory_insert(content)",
            "archival_memory_search(query, page)",
        ]),
        ("MemFS", WARN, [
            "memfs.create / read / write",
            "memfs.list / move / delete / grep",
            "memfs.history / rollback",
        ]),
        ("Comunicación", HOT, [
            "send_message(text)",
        ]),
    ]
    x0 = Inches(0.55)
    cw = Inches(3.05)
    for i, (title, c, items) in enumerate(groups):
        x = x0 + (cw + Inches(0.1)) * i
        rect(s, x, Inches(2.1), cw, Inches(4.7), fill=SURFACE, rounded=True)
        rect(s, x, Inches(2.1), cw, Inches(0.6), fill=c)
        textbox(s, x + Inches(0.25), Inches(2.2), cw - Inches(0.3), Inches(0.45),
                title, size=14, bold=True, color=BG, font="Georgia")
        for j, item in enumerate(items):
            y = Inches(2.95 + j * 0.6)
            rect(s, x + Inches(0.2), y, cw - Inches(0.4), Inches(0.5),
                 fill=BG, rounded=True)
            textbox(s, x + Inches(0.35), y + Inches(0.1), cw - Inches(0.7), Inches(0.35),
                    item, size=11, color=INK, font="Consolas")
    footer(s, n, total)


def slide_design_choices(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Decisiones de diseño", "Lo que aprendí construyéndolo")

    items = [
        ("Bloques atómicos al expulsar",
         "Nunca separar un AIMessage con tool_calls de sus ToolMessages.\n"
         "Evita los bugs #111/#126 de langmem.",
         ACCENT),
        ("Recursive summary fuera de messages",
         "Vive en su propio campo del estado y se inyecta como SystemMessage.\n"
         "El reducer append-only no admite slots mutables.",
         ACCENT_2),
        ("parallel_tool_calls = False",
         "Dos tools a la vez chocan en core_memory (un solo campo del state).\n"
         "MemGPT está pensado para chaining, no fan-out.",
         WARN),
        ("Síncrono a propósito",
         "El grafo se invoca con .invoke(); mezclar nodos async\n"
         "rompe con TypeError en LangGraph.",
         HOT),
    ]
    for i, (t, body, c) in enumerate(items):
        col, row = i % 2, i // 2
        x = Inches(0.55) + Inches(6.25) * col
        y = Inches(2.05) + Inches(2.5) * row
        rect(s, x, y, Inches(6.0), Inches(2.3), fill=SURFACE, rounded=True)
        rect(s, x, y, Inches(0.16), Inches(2.3), fill=c)
        textbox(s, x + Inches(0.35), y + Inches(0.2), Inches(5.5), Inches(0.5),
                t, size=15, bold=True, color=c, font="Georgia")
        textbox(s, x + Inches(0.35), y + Inches(0.75), Inches(5.5), Inches(1.5),
                body, size=12, color=INK, line_spacing=1.4)
    footer(s, n, total)


def slide_summary(prs, n, total):
    s = add_blank(prs)
    fill_bg(s)
    header(s, "Resumen", "Qué te llevas de MemGPT")

    points = [
        ("LLM como CPU", "La ventana es RAM, lo externo es disco.", ACCENT),
        ("Memoria jerárquica",
         "Working Context · FIFO · Recall · Archival · MemFS.",
         ACCENT_2),
        ("Auto-gestión",
         "El propio modelo decide qué guardar, qué buscar y cuántas veces.",
         WARN),
        ("Control de flujo",
         "Queue Manager + umbrales (70 % avisa, 100 % expulsa y resume).",
         HOT),
        ("Eventos automáticos",
         "Wall-clock + sleep-time agents para consolidar sin usuario.",
         ACCENT),
    ]
    for i, (k, v, c) in enumerate(points):
        y = Inches(2.15 + i * 0.92)
        rect(s, Inches(0.55), y, Inches(12.3), Inches(0.78), fill=SURFACE, rounded=True)
        rect(s, Inches(0.55), y, Inches(0.16), Inches(0.78), fill=c)
        textbox(s, Inches(0.85), y + Inches(0.18), Inches(3.5), Inches(0.5),
                k, size=18, bold=True, color=c, font="Georgia")
        textbox(s, Inches(4.4), y + Inches(0.2), Inches(8.4), Inches(0.5),
                v, size=14, color=INK)
    footer(s, n, total)


def slide_thanks(prs, n, total):
    s = add_blank(prs)
    fill_bg(s, BG)
    # decoración (espejo de la portada)
    band = rect(s, Inches(-2), Inches(-1), Inches(7), Inches(10), fill=SURFACE)
    band.rotation = -12
    band2 = rect(s, Inches(-1), Inches(0.5), Inches(6), Inches(8), fill=SURFACE_2)
    band2.rotation = -12
    slide_oval(s, Inches(0.3), Inches(1.4), Inches(2.6), ACCENT)
    slide_oval(s, Inches(2.0), Inches(4.8), Inches(1.6), ACCENT_2)
    slide_oval(s, Inches(-0.5), Inches(5.2), Inches(1.4), HOT)

    textbox(s, Inches(5.0), Inches(1.8), Inches(8), Inches(1.5),
            "Gracias", size=96, bold=True, color=INK, font="Georgia")
    textbox(s, Inches(5.0), Inches(3.4), Inches(8), Inches(0.6),
            "Preguntas, comentarios, ideas.", size=24, color=DIM)

    # Recursos
    rect(s, Inches(5.0), Inches(4.5), Inches(7.5), Inches(2.3), fill=SURFACE, rounded=True)
    textbox(s, Inches(5.25), Inches(4.65), Inches(7.0), Inches(0.5),
            "Recursos", size=16, bold=True, color=ACCENT, font="Georgia")
    res = [
        ("Paper", "arXiv:2310.08560", ACCENT),
        ("Letta (sucesor)", "github.com/letta-ai/letta", ACCENT_2),
        ("Mi implementación", "github.com/maximofn/memGPT", WARN),
        ("Blog", "maximofn.com", HOT),
    ]
    for i, (k, v, c) in enumerate(res):
        y = Inches(5.15 + i * 0.38)
        dot(s, Inches(5.25), y + Inches(0.08), color=c, size=Inches(0.12))
        textbox(s, Inches(5.5), y, Inches(2.5), Inches(0.35),
                k, size=12, bold=True, color=INK)
        textbox(s, Inches(8.1), y, Inches(4.3), Inches(0.35),
                v, size=12, color=DIM, font="Consolas")


# ---------- Build ----------
def main():
    prs = new_deck()
    slides = [
        slide_title,
        slide_problem,
        slide_os_analogy,
        slide_three_pillars,
        slide_context_window,
        slide_working_context,
        slide_core_memory_tools,
        slide_fifo,
        slide_thresholds,
        slide_recursive_summary,
        slide_recall,
        slide_archival,
        slide_archival_vs_rag,
        slide_memfs,
        slide_compare_layers,
        slide_decision_matrix,
        slide_cycle,
        slide_events,
        slide_heartbeat,
        slide_my_impl_overview,
        slide_my_tools,
        slide_design_choices,
        slide_summary,
        slide_thanks,
    ]
    total = len(slides)
    # title and thanks no footer
    slides[0](prs)
    for i, fn in enumerate(slides[1:-1], start=2):
        fn(prs, i, total)
    slides[-1](prs, total, total)
    prs.save(OUT)
    print(f"✓ Saved: {OUT}  ({total} slides)")


if __name__ == "__main__":
    main()
